"""
PowerPoint (.pptx) converter for Human-to-AI Format Converter.
Extracts text from slides including titles, content, and notes.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from .base import BaseConverter


class PptxConverter(BaseConverter):
    """Converter for Microsoft PowerPoint presentations."""
    
    def _get_file_type(self) -> str:
        return "powerpoint_presentation"
    
    def _extract_content(self) -> dict:
        """Extract content from PowerPoint presentation."""
        prs = Presentation(self.file_path)
        
        content = {
            "text": "",
            "slides": [],
            "total_slides": len(prs.slides)
        }
        
        full_text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "notes": "",
                "shapes": []
            }
            
            # Extract title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text.strip()
                full_text_parts.append(slide_data["title"])
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            shape_text.append(para_text)
                    
                    if shape_text:
                        # Don't duplicate the title
                        if shape != slide.shapes.title:
                            slide_data["content"].extend(shape_text)
                            full_text_parts.extend(shape_text)
                        
                        slide_data["shapes"].append({
                            "type": shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
                            "text": shape_text
                        })
                
                # Handle tables in shapes
                if shape.has_table:
                    table = shape.table
                    table_data = {
                        "rows": []
                    }
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data["rows"].append(row_data)
                    
                    if table_data["rows"]:
                        table_data["headers"] = table_data["rows"][0] if table_data["rows"] else []
                        slide_data["shapes"].append({
                            "type": "TABLE",
                            "table": table_data
                        })
            
            # Extract notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    slide_data["notes"] = notes_frame.text.strip()
            
            content["slides"].append(slide_data)
        
        content["text"] = "\n\n".join(full_text_parts)
        
        # Update metadata
        self._metadata = {
            "slide_count": len(prs.slides),
            "word_count": len(content["text"].split())
        }
        
        return content
    
    def _format_as_markdown(self, content: dict) -> str:
        """Format PowerPoint content as Markdown."""
        md_parts = []
        
        for slide in content["slides"]:
            # Slide header
            md_parts.append(f"## Slide {slide['slide_number']}")
            
            # Title
            if slide["title"]:
                md_parts.append(f"### {slide['title']}")
            
            md_parts.append("")
            
            # Content
            for text in slide["content"]:
                # Format as bullet points
                md_parts.append(f"- {text}")
            
            # Tables from shapes
            for shape in slide["shapes"]:
                if shape["type"] == "TABLE" and "table" in shape:
                    table = shape["table"]
                    if table["rows"]:
                        md_parts.append("")
                        
                        # Header row
                        if table.get("headers"):
                            md_parts.append("| " + " | ".join(table["headers"]) + " |")
                            md_parts.append("| " + " | ".join(["---"] * len(table["headers"])) + " |")
                        
                        # Data rows
                        start_idx = 1 if table.get("headers") else 0
                        for row in table["rows"][start_idx:]:
                            md_parts.append("| " + " | ".join(row) + " |")
                        
                        md_parts.append("")
            
            # Notes
            if slide["notes"]:
                md_parts.append("")
                md_parts.append(f"**Speaker Notes:** {slide['notes']}")
            
            md_parts.append("")
            md_parts.append("---")
            md_parts.append("")
        
        return "\n".join(md_parts)
